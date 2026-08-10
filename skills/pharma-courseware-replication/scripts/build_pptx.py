#!/usr/bin/env python3
"""Build an editable 16:9 PPTX from courseware.content.json + style tokens.

Usage:
  python3 build_pptx.py <content.json> <out.pptx> [--tokens tokens.json]
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _rgb(hex6: str) -> RGBColor:
    h = hex6.strip().lstrip("#")
    if len(h) == 3:
        h = "".join(ch * 2 for ch in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _font(run, size_pt: float, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "PingFang SC"


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    _font(run, size, color, bold)
    return box


def add_lines(slide, left, top, width, height, lines, size, color, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(size * 0.35)
        run = p.add_run()
        run.text = str(line)
        _font(run, size, color, bold)
    return box


def fill_shape(shape, hex6: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex6)
    shape.line.fill.background()


def stroke_shape(shape, hex6: str, pt: float = 1.25) -> None:
    shape.line.color.rgb = _rgb(hex6)
    shape.line.width = Pt(pt)


def card(slide, left, top, width, height, colors):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, colors["surface"])
    stroke_shape(shape, colors["line"])
    return shape


def header_bar(slide, colors, W) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.12))
    fill_shape(bar, colors["primary"])


def footer(slide, meta, page_idx, page_count, colors, fonts, W, H, ml, mr) -> None:
    y = H - Inches(0.38)
    add_text(
        slide,
        ml,
        y,
        Inches(4),
        Inches(0.3),
        meta.get("organization") or "",
        fonts["footer"],
        _rgb(colors["ink_muted"]),
    )
    add_text(
        slide,
        Inches(4.5),
        y,
        Inches(4.5),
        Inches(0.3),
        meta.get("internal_notice") or "仅限于内部学习",
        fonts["footer"],
        _rgb(colors["ink_muted"]),
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        W - mr - Inches(1.2),
        y,
        Inches(1.2),
        Inches(0.3),
        f"{page_idx:02d} / {page_count:02d}",
        fonts["footer"],
        _rgb(colors["ink_muted"]),
        align=PP_ALIGN.RIGHT,
    )


def page_title(slide, title, colors, fonts, ml, mt) -> None:
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, ml, mt, Inches(0.12), Inches(0.38)
    )
    fill_shape(chip, colors["primary"])
    add_text(
        slide,
        ml + Inches(0.28),
        mt - Inches(0.02),
        Inches(11),
        Inches(0.45),
        title or "",
        fonts["page_title"],
        _rgb(colors["ink"]),
        bold=True,
    )


def render_cover(slide, s, meta, colors, fonts, W, H, ml) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), H)
    fill_shape(panel, colors["primary"])
    soft = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0), Inches(0.15), H
    )
    fill_shape(soft, colors["primary_soft"])
    add_text(
        slide,
        ml + Inches(0.3),
        Inches(2.1),
        Inches(11),
        Inches(1.2),
        s.get("title") or meta.get("title") or "",
        fonts["cover_title"],
        _rgb(colors["ink"]),
        bold=True,
    )
    add_text(
        slide,
        ml + Inches(0.3),
        Inches(3.4),
        Inches(10),
        Inches(0.5),
        s.get("subtitle") or meta.get("subtitle") or "",
        fonts["body"] + 2,
        _rgb(colors["ink_muted"]),
    )
    if s.get("body"):
        add_text(
            slide,
            ml + Inches(0.3),
            Inches(4.2),
            Inches(10),
            Inches(1.2),
            s.get("body"),
            fonts["body"],
            _rgb(colors["ink"]),
        )
    add_text(
        slide,
        ml + Inches(0.3),
        H - Inches(0.5),
        Inches(10),
        Inches(0.3),
        meta.get("internal_notice") or "仅限于内部学习",
        fonts["footer"],
        _rgb(colors["ink_muted"]),
    )


def render_agenda(slide, s, colors, fonts, ml, mt, content_w) -> None:
    page_title(slide, s.get("title") or "目录", colors, fonts, ml, mt)
    items = s.get("bullets") or s.get("items") or []
    y = mt + Inches(0.7)
    for i, item in enumerate(items):
        card(slide, ml, y, content_w, Inches(0.72), colors)
        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, ml + Inches(0.18), y + Inches(0.16), Inches(0.4), Inches(0.4)
        )
        fill_shape(num, colors["primary"])
        add_text(
            slide,
            ml + Inches(0.18),
            y + Inches(0.2),
            Inches(0.4),
            Inches(0.35),
            f"{i + 1:02d}",
            fonts["chip"],
            _rgb(colors["white"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text = item if isinstance(item, str) else str(item.get("title") or item)
        add_text(
            slide,
            ml + Inches(0.75),
            y + Inches(0.2),
            content_w - Inches(1),
            Inches(0.4),
            text,
            fonts["card_title"],
            _rgb(colors["ink"]),
            bold=True,
        )
        y += Inches(0.85)


def render_cards(slide, s, colors, fonts, ml, mt, content_w) -> None:
    page_title(slide, s.get("title") or "", colors, fonts, ml, mt)
    cards = list(s.get("cards") or [])
    bullets = s.get("bullets") or []
    if not cards and bullets:
        cards = [{"title": b, "body": ""} for b in bullets]
    n = max(len(cards), 1)
    cols = 3 if n == 3 or n > 4 else 2
    gap = Inches(0.2)
    cw = (content_w - gap * (cols - 1)) / cols
    ch = Inches(2.15)
    y0 = mt + Inches(0.75)
    for i, cdata in enumerate(cards):
        col = i % cols
        row = i // cols
        left = ml + col * (cw + gap)
        top = y0 + row * (ch + gap)
        card(slide, left, top, cw, ch, colors)
        title = cdata.get("title") if isinstance(cdata, dict) else str(cdata)
        body = cdata.get("body", "") if isinstance(cdata, dict) else ""
        add_text(
            slide,
            left + Inches(0.18),
            top + Inches(0.18),
            cw - Inches(0.36),
            Inches(0.45),
            title or "",
            fonts["card_title"],
            _rgb(colors["primary"]),
            bold=True,
        )
        if body:
            add_text(
                slide,
                left + Inches(0.18),
                top + Inches(0.65),
                cw - Inches(0.36),
                ch - Inches(0.85),
                body,
                fonts["body"],
                _rgb(colors["ink"]),
            )


def render_body(slide, s, colors, fonts, ml, mt, content_w, H, mb) -> None:
    page_title(slide, s.get("title") or "", colors, fonts, ml, mt)
    y = mt + Inches(0.7)
    if s.get("body"):
        add_text(
            slide,
            ml,
            y,
            content_w,
            Inches(1.6),
            s.get("body"),
            fonts["body"],
            _rgb(colors["ink"]),
        )
        y += Inches(1.5)
    if s.get("bullets"):
        lines = [f"• {b}" for b in s["bullets"]]
        add_lines(
            slide,
            ml,
            y,
            content_w,
            H - y - mb - Inches(0.5),
            lines,
            fonts["body"],
            _rgb(colors["ink"]),
        )


def render_table(slide, s, colors, fonts, ml, mt, content_w) -> None:
    page_title(slide, s.get("title") or "", colors, fonts, ml, mt)
    table = s.get("table") or {}
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers and not rows:
        render_body(slide, s, colors, fonts, ml, mt, content_w, Inches(7.5), Inches(0.4))
        return
    cols = max(len(headers), max((len(r) for r in rows), default=1))
    if len(headers) < cols:
        headers = list(headers) + [""] * (cols - len(headers))
    n_rows = 1 + len(rows)
    left, top = ml, mt + Inches(0.75)
    height = Inches(min(5.2, 0.45 * n_rows + 0.4))
    shape = slide.shapes.add_table(n_rows, cols, left, top, content_w, height)
    tbl = shape.table
    for c in range(cols):
        cell = tbl.cell(0, c)
        cell.text = headers[c] if c < len(headers) else ""
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(fonts["body"])
                r.font.color.rgb = _rgb(colors["white"])
                r.font.name = "PingFang SC"
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(colors["header_bar"])
    for ri, row in enumerate(rows):
        for c in range(cols):
            cell = tbl.cell(ri + 1, c)
            cell.text = str(row[c]) if c < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(max(fonts["body"] - 1, 10))
                    r.font.color.rgb = _rgb(colors["ink"])
                    r.font.name = "PingFang SC"
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                colors["surface"] if ri % 2 == 0 else colors["primary_soft"]
            )


def render_product(slide, s, colors, fonts, ml, mt, content_w) -> None:
    page_title(slide, s.get("title") or "", colors, fonts, ml, mt)
    slot_w, slot_h = Inches(3.6), Inches(4.2)
    card(slide, ml, mt + Inches(0.75), slot_w, slot_h, colors)
    add_text(
        slide,
        ml + Inches(0.3),
        mt + Inches(2.4),
        slot_w - Inches(0.6),
        Inches(0.8),
        "【包装图位】\n业务授权实拍",
        fonts["body"],
        _rgb(colors["ink_muted"]),
        align=PP_ALIGN.CENTER,
    )
    right_l = ml + slot_w + Inches(0.3)
    right_w = content_w - slot_w - Inches(0.3)
    if s.get("body"):
        add_text(
            slide,
            right_l,
            mt + Inches(0.75),
            right_w,
            Inches(1.2),
            s.get("body"),
            fonts["body"],
            _rgb(colors["ink"]),
        )
    if s.get("bullets"):
        add_lines(
            slide,
            right_l,
            mt + Inches(2.0),
            right_w,
            Inches(3.5),
            [f"• {b}" for b in s["bullets"]],
            fonts["body"],
            _rgb(colors["ink"]),
        )


def render_slide(prs, s, meta, tokens, page_idx, page_count) -> None:
    colors = tokens["colors"]
    fonts = tokens["fonts"]["sizes_pt"]
    lay = tokens["layout"]
    W = Inches(tokens["slide"]["width_in"])
    H = Inches(tokens["slide"]["height_in"])
    ml = Inches(lay["margin_left_in"])
    mr = Inches(lay["margin_right_in"])
    mt = Inches(lay["margin_top_in"])
    mb = Inches(lay["margin_bottom_in"])
    content_w = W - ml - mr

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    fill_shape(bg, colors["bg"])

    ptype = (s.get("page_type") or "custom").lower()
    header_bar(slide, colors, W)

    if ptype == "cover":
        render_cover(slide, s, meta, colors, fonts, W, H, ml)
    elif ptype == "agenda":
        render_agenda(slide, s, colors, fonts, ml, mt, content_w)
    elif ptype in {"product_intro", "ingredient"}:
        render_product(slide, s, colors, fonts, ml, mt, content_w)
    elif ptype in {"comparison", "usage", "caution"} and s.get("table"):
        render_table(slide, s, colors, fonts, ml, mt, content_w)
    elif ptype in {"symptoms", "summary", "indication_scene", "objectives"} or s.get("cards"):
        render_cards(slide, s, colors, fonts, ml, mt, content_w)
    else:
        if s.get("table"):
            render_table(slide, s, colors, fonts, ml, mt, content_w)
        elif s.get("cards"):
            render_cards(slide, s, colors, fonts, ml, mt, content_w)
        else:
            render_body(slide, s, colors, fonts, ml, mt, content_w, H, mb)

    if ptype != "cover":
        footer(slide, meta, page_idx, page_count, colors, fonts, W, H, ml, mr)


def build(content_path: Path, out_path: Path, tokens_path: Path | None) -> None:
    data = json.loads(content_path.read_text(encoding="utf-8"))
    if tokens_path is None:
        # Provisional neutral chrome for structure smoke tests only.
        # Real deposits must pass tokens extracted from the business reference PPT.
        # Do NOT use store-vitality-v1 here (that id is illustration-only).
        tokens_path = (
            Path(__file__).resolve().parent.parent
            / "references"
            / "styles"
            / "pptx-shell-neutral-v1"
            / "tokens.json"
        )
    tokens = json.loads(tokens_path.read_text(encoding="utf-8"))
    meta = data.get("meta") or {}
    slides = data.get("slides") or []
    if not slides:
        raise SystemExit("content.json has no slides[]")

    prs = Presentation()
    prs.slide_width = Inches(tokens["slide"]["width_in"])
    prs.slide_height = Inches(tokens["slide"]["height_in"])

    n = len(slides)
    for i, s in enumerate(slides, 1):
        render_slide(prs, s, meta, tokens, i, n)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path} ({n} slides)")


def main(argv=None) -> None:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("content_json", type=Path)
    p.add_argument("out_pptx", type=Path)
    p.add_argument("--tokens", type=Path, default=None)
    args = p.parse_args(argv)
    build(args.content_json, args.out_pptx, args.tokens)


if __name__ == "__main__":
    main()
