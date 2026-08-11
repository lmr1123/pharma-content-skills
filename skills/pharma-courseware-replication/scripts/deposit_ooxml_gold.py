#!/usr/bin/env python3
"""
OOXML 金样归档（近 100% 视觉复刻的唯一诚实路径）

用法：
  python3 scripts/deposit_ooxml_gold.py \\
    --source "/path/to/参考.pptx" \\
    --template-id kangaisen-lycopene-health-edu-v1 \\
    [--name-zh "番茄红素健康科普"] \\
    [--open]

做什么：
  1. 把参考 PPTX **原样拷贝**为模板样片（视觉 = 100% 原片，不是 pptxgenjs 重画）
  2. 写 inventory.json（页数 / 媒体数 / 体积 / SHA256）
  3. 写 template-manifest.md（fidelity: gold-aligned-ooxml-v1）
  4. 写 content 槽位粗清单（按 shape 文本顺序，便于换主题）
  5. 可选 open 样片

禁止：
  - 用本脚本却把 fidelity 标成 path-only 的框架引擎
  - 用通用壳 / 形状重画冒充 gold-aligned

换主题近 100%：
  保留本金样 OOXML，只换文字槽与图槽（生产仓 ingredient-health-edu 用 artifact-tool
  按 shape id 换槽）。Skill 内尚未自包含 artifact-tool 时：金样交付 = 本归档；
  量产换题走生产引擎或后续迁入的 ooxml 换槽器。
"""

from __future__ import annotations

import argparse
import hashlib
import json
import shutil
import subprocess
import sys
import zipfile
from datetime import date
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: need python-pptx  (pip install python-pptx)", file=sys.stderr)
    sys.exit(2)

ROOT = Path(__file__).resolve().parents[1]  # skills/pharma-courseware-replication
REPO = ROOT.parents[1] if (ROOT.parents[1] / "workspace").exists() else ROOT
# Prefer repo workspace (git clone); fall back to skill-local workspace (WorkBuddy install layout)
for candidate in (REPO / "workspace", ROOT / "workspace"):
    if candidate.parent.exists():
        WORKSPACE = candidate
        break
else:
    WORKSPACE = ROOT / "workspace"


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def inventory_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    slides_meta = []
    text_slots = []
    for i, slide in enumerate(prs.slides, start=1):
        shape_count = 0
        page_texts = []
        for shape in slide.shapes:
            shape_count += 1
            if not getattr(shape, "has_text_frame", False):
                continue
            raw = (shape.text_frame.text or "").strip()
            if not raw:
                continue
            slot = {
                "slide": i,
                "shape_name": getattr(shape, "name", "") or "",
                "text": raw[:500],
                "chars": len(raw),
            }
            page_texts.append(slot)
            text_slots.append(slot)
        slides_meta.append({"slide": i, "shapes": shape_count, "text_shapes": len(page_texts)})

    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        media = [n for n in names if n.startswith("ppt/media/") and not n.endswith("/")]

    return {
        "slides": len(prs.slides),
        "slide_width_emu": int(prs.slide_width),
        "slide_height_emu": int(prs.slide_height),
        "shapes_total": sum(s["shapes"] for s in slides_meta),
        "media_files": len(media),
        "media_ext_counts": _ext_counts(media),
        "size_bytes": path.stat().st_size,
        "sha256": sha256_file(path),
        "slides_meta": slides_meta,
        "text_slot_count": len(text_slots),
        "text_slots_preview": text_slots[:40],
        "text_slots": text_slots,
    }


def _ext_counts(media: list[str]) -> dict:
    counts: dict[str, int] = {}
    for n in media:
        ext = Path(n).suffix.lower() or "(none)"
        counts[ext] = counts.get(ext, 0) + 1
    return dict(sorted(counts.items(), key=lambda x: (-x[1], x[0])))


def rich_design_gate(inv: dict) -> list[str]:
    """If any fire, pptxgenjs framework rebuild is FORBIDDEN as gold-aligned."""
    reasons = []
    if inv["media_files"] >= 15:
        reasons.append(f"media_files={inv['media_files']} (≥15 → 必须 OOXML 原片，禁止形状重画冒充金样)")
    if inv["size_bytes"] >= 2_000_000:
        reasons.append(f"size={inv['size_bytes']} (≥2MB → 通常含纹理/真图)")
    if any(ext in inv["media_ext_counts"] for ext in (".svg", ".emf", ".wmf")):
        reasons.append("含 SVG/矢量图标 → 框架壳无法保留")
    return reasons


def write_manifest(
    path: Path,
    *,
    template_id: str,
    name_zh: str,
    source: Path,
    inv: dict,
    gate: list[str],
) -> None:
    path.write_text(
        f"""# 模板清单 · {template_id}

> 业务自有模板（`workspace/` · `git pull` 不覆盖）

## 一句话

**OOXML 金样归档**：样片 = 参考原片本体（近 100% 视觉）。**不是** pptxgenjs 框架重画。

## 真源

| 项 | 值 |
|----|-----|
| 模板 id | `{template_id}` |
| 中文名 | {name_zh} |
| 参考源（本机） | `{source}` |
| 样片 | `output/courseware.pptx`（= 原片拷贝） |
| 页数 | {inv["slides"]} |
| 媒体文件 | {inv["media_files"]} |
| 文本槽（粗） | {inv["text_slot_count"]} |
| 体积 | {inv["size_bytes"]} bytes |
| SHA256 | `{inv["sha256"]}` |
| 归档日 | {date.today().isoformat()} |

## 保真声明

```text
fidelity: gold-aligned-ooxml-v1
strategy: ooxml-gold-archive
engine: scripts/deposit_ooxml_gold.py  （金样样片）
theme_extension: ooxml-slot-replace（生产 ingredient-health-edu 或后续迁入换槽器）
```

### 为何不能标 framework / health-popularization pptxgenjs 为 gold-aligned

触发富设计门禁：

{chr(10).join("- " + r for r in gate) or "- （无）"}

框架壳只有页骨架与近似色，**没有**原片媒体/阴影/自由曲线/嵌入字体观感。  
禁止把 `fidelity: gold-aligned-v1` 写在 pptxgenjs 重画产物上。

## 出片 / 换主题

### 金样演示（业务看「复刻效果」）

直接打开 `output/courseware.pptx` —— 与参考原片一致。

### 换新主题（量产）

1. 复制本金样 OOXML，**不要**用通用壳重画  
2. 按 `inventory.json` 的 text_slots / 生产 107 文字槽 + 69 图槽换内容与授权图  
3. 生产仓路径（完整换槽）：`ingredient-health-edu-pptx-v1` + `ingredient-health-edu-theme/v1`  
4. Skill 内若仅有归档：先交付金样观感；换题引擎迁入前不得声称「换题也 100%」

## 目录

```text
{template_id}/
  template-manifest.md
  inventory.json
  reuse/change-list.md
  output/courseware.pptx      # 100% 原片
  samples/source-path.txt
```
""",
        encoding="utf-8",
    )


def write_change_list(path: Path, inv: dict) -> None:
    path.write_text(
        f"""# 换主题改什么 / 不改什么

## 不改（锁金样）

- 画布尺寸、母版、版式布局几何
- 背景纹理、图标 SVG、阴影、自由曲线封面
- 页序与页型节奏（当前 {inv["slides"]} 页）
- 非本主题的装饰 chrome

## 改

- 各页文案（见 `inventory.json` → `text_slots`）
- 主题相关配图（原图槽换业务授权/新生成图；禁止假包装）
- 页脚机构名等业务字段

## 禁止

- 用 pptxgenjs / 通用圆角卡「重画一版」当金样
- 打开金样只改字却不换图就当新主题正式交付（医学/品牌串货）
- 把 `fidelity` 标成 `gold-aligned` 却交付框架壳
""",
        encoding="utf-8",
    )


def main() -> int:
    ap = argparse.ArgumentParser(description="Archive reference PPTX as OOXML gold template")
    ap.add_argument("--source", required=True, help="参考原片 .pptx 路径")
    ap.add_argument("--template-id", required=True, help="写入 workspace/templates/<id>/")
    ap.add_argument("--name-zh", default="", help="中文展示名")
    ap.add_argument("--open", action="store_true", help="完成后 open 样片")
    ap.add_argument(
        "--workspace",
        default=str(WORKSPACE),
        help="workspace 根（默认自动探测仓库或 Skill 内 workspace）",
    )
    args = ap.parse_args()

    source = Path(args.source).expanduser().resolve()
    if not source.is_file():
        print(f"ERROR: source not found: {source}", file=sys.stderr)
        return 2
    if source.suffix.lower() != ".pptx":
        print("ERROR: source must be .pptx", file=sys.stderr)
        return 2

    template_id = args.template_id.strip().replace(" ", "-")
    name_zh = args.name_zh.strip() or template_id
    ws = Path(args.workspace).expanduser().resolve()
    dest_root = ws / "templates" / template_id
    out_dir = dest_root / "output"
    reuse_dir = dest_root / "reuse"
    samples_dir = dest_root / "samples"
    for d in (out_dir, reuse_dir, samples_dir):
        d.mkdir(parents=True, exist_ok=True)

    inv = inventory_pptx(source)
    gate = rich_design_gate(inv)
    if not gate:
        gate = ["媒体较少也可能是设计精稿；默认仍以原片归档为金样，禁止未经验收的框架壳 gold-aligned"]

    gold_out = out_dir / "courseware.pptx"
    shutil.copy2(source, gold_out)
    # 再读拷贝校验
    inv_out = inventory_pptx(gold_out)
    if inv_out["sha256"] != inv["sha256"]:
        print("ERROR: copy sha256 mismatch", file=sys.stderr)
        return 1

    (dest_root / "inventory.json").write_text(
        json.dumps(
            {
                "template_id": template_id,
                "name_zh": name_zh,
                "source": str(source),
                "fidelity": "gold-aligned-ooxml-v1",
                "strategy": "ooxml-gold-archive",
                "rich_design_gate": gate,
                "inventory": inv_out,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    (samples_dir / "source-path.txt").write_text(str(source) + "\n", encoding="utf-8")
    write_manifest(
        dest_root / "template-manifest.md",
        template_id=template_id,
        name_zh=name_zh,
        source=source,
        inv=inv_out,
        gate=gate,
    )
    write_change_list(reuse_dir / "change-list.md", inv_out)

    # 内容初稿粗稿：按页列出原文文本，便于业务复核
    draft_lines = [
        f"# 内容粗抽 · {name_zh}",
        "",
        "> 来源：参考原片文本槽（顺序抽取）。换主题时替换对应段落；**版式以 OOXML 金样为准**。",
        f"> 样片已归档：`output/courseware.pptx`（SHA256 `{inv_out['sha256'][:16]}…`）",
        "",
    ]
    by_slide: dict[int, list[str]] = {}
    for slot in inv_out["text_slots"]:
        by_slide.setdefault(slot["slide"], []).append(slot["text"])
    for s in range(1, inv_out["slides"] + 1):
        draft_lines.append(f"## 第 {s} 页")
        for t in by_slide.get(s, ["（本页无文本框或仅图片）"]):
            draft_lines.append(f"- {t.replace(chr(10), ' / ')}")
        draft_lines.append("")
    draft_path = reuse_dir / "content-draft.md"
    draft_path.write_text("\n".join(draft_lines), encoding="utf-8")

    print(
        json.dumps(
            {
                "ok": True,
                "template_id": template_id,
                "dest": str(dest_root),
                "pptx": str(gold_out),
                "fidelity": "gold-aligned-ooxml-v1",
                "slides": inv_out["slides"],
                "media_files": inv_out["media_files"],
                "size_bytes": inv_out["size_bytes"],
                "sha256": inv_out["sha256"],
                "rich_design_gate": gate,
                "content_draft": str(draft_path),
            },
            ensure_ascii=False,
            indent=2,
        )
    )

    if args.open:
        for p in (gold_out, draft_path):
            try:
                subprocess.run(["open", str(p)], check=False)
            except Exception:
                pass

    return 0


if __name__ == "__main__":
    sys.exit(main())
